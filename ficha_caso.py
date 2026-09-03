"""
Generador de la "Ficha de Caso" individual (7 slides, pptx independiente
del Reporte de Cartera):

1. Resumen del caso (Ficha del Caso).
2. Descripción del Siniestro / Materia Asegurada (texto libre, en blanco).
3. Registro Fotográfico: si se le pasan fotos ya seleccionadas (p.ej.
   extraídas de un Acta de Inspección con acta_inspeccion.py), se generan
   tantas páginas de 6 fotos (con su pie de foto) como hagan falta; si no,
   una sola slide con 6 casilleros en blanco para completar a mano.
4. Reserva del caso: valor y desglose extraídos de la planilla base
   (Pérdida bruta, Deducible, Monto asegurado, Pérdida neta/Reserva,
   Gastos, Honorarios), con una columna en blanco para que el ajustador
   justifique cada concepto.
5. Estado Actual del Siniestro: precargado con la observación completa
   del Pipeline para ese caso (si existe), editable; si no hay dato,
   queda en blanco como texto libre.
6. Detalle de Gestiones Realizadas (tabla de 6 filas vacías, Fecha +
   Detalle, a llenar a mano).
7. Próximas Acciones (texto libre, en blanco).

Las slides de texto libre en blanco (2 y 7, y la 5 cuando no hay
observación del Pipeline) llevan la nota "(Texto libre)".

No usa una plantilla .pptx: arma las slides desde cero con python-pptx,
reutilizando el mismo logo y paleta de colores (navy + teal) que
assets/plantilla_estado_cartera.pptx.
"""
import io
import os

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Inches, Pt

LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "logo_jpv.png")

SLIDE_WIDTH = Emu(9144000)
SLIDE_HEIGHT = Emu(5143500)

NAVY_OSCURO = RGBColor(0x0D, 0x1F, 0x38)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x14, 0xA8, 0xA0)
TEAL_OSCURO = RGBColor(0x0D, 0x73, 0x77)
GRIS_TEXTO = RGBColor(0x33, 0x33, 0x33)
GRIS_CLARO = RGBColor(0xF4, 0xF6, 0xF8)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

HEADER_ALTO = Inches(0.82)


def _truncar(texto, max_chars):
    texto = str(texto).strip()
    if len(texto) <= max_chars:
        return texto
    return texto[: max_chars - 1].rstrip() + "…"


def _sin_relleno_ni_borde(shape):
    shape.fill.background()
    shape.line.fill.background()


def _agregar_header(slide, titulo, subtitulo):
    barra = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, HEADER_ALTO)
    barra.fill.solid()
    barra.fill.fore_color.rgb = NAVY
    barra.line.fill.background()
    barra.shadow.inherit = False

    divisor = slide.shapes.add_shape(1, 0, HEADER_ALTO, SLIDE_WIDTH, Emu(36576))
    divisor.fill.solid()
    divisor.fill.fore_color.rgb = TEAL_OSCURO
    divisor.line.fill.background()
    divisor.shadow.inherit = False

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Emu(48142), Emu(76341), height=Emu(528226))

    caja_titulo = slide.shapes.add_textbox(Inches(1.25), Inches(0.12), Inches(7.5), Inches(0.6))
    tf = caja_titulo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, BLANCO

    p2 = tf.add_paragraph()
    p2.text = subtitulo
    p2.font.size, p2.font.italic, p2.font.color.rgb = Pt(12), True, TEAL


def _campo(slide, left, top, width, etiqueta, valor):
    caja = slide.shapes.add_textbox(left, top, width, Inches(0.62))
    tf = caja.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = etiqueta.upper()
    p.font.size, p.font.bold, p.font.color.rgb = Pt(9), True, TEAL_OSCURO

    p2 = tf.add_paragraph()
    p2.text = str(valor) if valor not in (None, "", "nan", "NaT") else "—"
    p2.font.size, p2.font.color.rgb = Pt(14), GRIS_TEXTO


def _slide_resumen(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _agregar_header(slide, "FICHA DE CASO", datos.get("nickname") or datos.get("asegurado") or "")

    col1_x, col2_x = Inches(0.5), Inches(5.15)
    ancho_col = Inches(4.3)
    top = Inches(1.15)
    paso = Inches(0.82)

    columna_1 = [
        ("Asegurado", datos.get("asegurado")),
        ("N° de Siniestro", datos.get("numero_siniestro")),
        ("Caso JPV", datos.get("caso_jpv")),
        ("Estado", datos.get("estado")),
        ("Monto asegurado", datos.get("monto_asegurado_fmt")),
    ]
    columna_2 = [
        ("Fecha de ocurrencia", datos.get("fecha_ocurrencia")),
        ("Fecha de denuncio", datos.get("fecha_denuncio")),
        ("Fecha de asignación", datos.get("fecha_asignacion")),
        ("Días desde asignación", datos.get("dias_asignacion")),
        ("Pérdida bruta", datos.get("perdida_bruta_fmt")),
    ]
    for i, (etiqueta, valor) in enumerate(columna_1):
        _campo(slide, col1_x, top + i * paso, ancho_col, etiqueta, valor)
    for i, (etiqueta, valor) in enumerate(columna_2):
        _campo(slide, col2_x, top + i * paso, ancho_col, etiqueta, valor)

    pie = slide.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(9), Inches(0.3))
    p = pie.text_frame.paragraphs[0]
    p.text = "JPV Asociados · Ajustadores Especializados"
    p.font.size, p.font.color.rgb = Pt(8), RGBColor(0x99, 0x99, 0x99)
    return slide


def _caja_manual(slide, left, top, width, height, titulo):
    """Etiqueta + caja vacía con borde punteado, pensada para completarse a
    mano en PowerPoint (mismo estilo que los casilleros de fotos)."""
    etiqueta = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = etiqueta.text_frame.paragraphs[0]
    p.text = titulo.upper()
    p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, TEAL_OSCURO

    alto_caja = height - Inches(0.32)
    caja = slide.shapes.add_shape(1, left, top + Inches(0.32), width, alto_caja)
    caja.fill.solid()
    caja.fill.fore_color.rgb = GRIS_CLARO
    caja.line.color.rgb = RGBColor(0xB0, 0xB8, 0xC2)
    caja.line.width = Pt(1)
    caja.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    caja.shadow.inherit = False
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = "(Texto libre)"
    p.font.size, p.font.italic, p.font.color.rgb = Pt(11), True, RGBColor(0x8A, 0x93, 0x9E)


def _slide_descripcion(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _agregar_header(slide, "DESCRIPCIÓN DEL SINIESTRO Y MATERIA ASEGURADA", datos.get("nickname") or datos.get("asegurado") or "")

    left, width = Inches(0.4), SLIDE_WIDTH - Inches(0.8)
    top = HEADER_ALTO + Emu(36576) + Inches(0.3)
    alto_caja, espacio = Inches(1.85), Inches(0.25)

    _caja_manual(slide, left, top, width, alto_caja, "Descripción del Siniestro")
    _caja_manual(slide, left, top + alto_caja + espacio, width, alto_caja, "Materia Asegurada")
    return slide


_FOTOS_POR_PAGINA = 6
_FOTOS_COLS, _FOTOS_FILAS = 3, 2


def _posiciones_grilla_fotos():
    margen, espacio = Inches(0.4), Inches(0.2)
    ancho_disponible = SLIDE_WIDTH - 2 * margen - (_FOTOS_COLS - 1) * espacio
    alto_disponible = SLIDE_HEIGHT - HEADER_ALTO - Emu(36576) - Inches(0.6) - (_FOTOS_FILAS - 1) * espacio
    ancho, alto = Emu(int(ancho_disponible / _FOTOS_COLS)), Emu(int(alto_disponible / _FOTOS_FILAS))
    top_inicial = HEADER_ALTO + Emu(36576) + Inches(0.3)
    posiciones = []
    for f in range(_FOTOS_FILAS):
        for c in range(_FOTOS_COLS):
            left = margen + c * (ancho + espacio)
            top = top_inicial + f * (alto + espacio)
            posiciones.append((left, top, ancho, alto))
    return posiciones


def _casillero_vacio(slide, left, top, ancho, alto, numero):
    marco = slide.shapes.add_shape(1, left, top, ancho, alto)
    marco.fill.solid()
    marco.fill.fore_color.rgb = GRIS_CLARO
    marco.line.color.rgb = RGBColor(0xB0, 0xB8, 0xC2)
    marco.line.width = Pt(1)
    marco.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    marco.shadow.inherit = False
    tf = marco.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"Foto {numero}"
    p.font.size, p.font.color.rgb = Pt(13), RGBColor(0x8A, 0x93, 0x9E)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _casillero_con_foto(slide, left, top, ancho, alto, imagen_bytes, pie):
    alto_pie = Inches(0.4)
    alto_imagen = alto - alto_pie

    marco = slide.shapes.add_shape(1, left, top, ancho, alto)
    marco.fill.solid()
    marco.fill.fore_color.rgb = BLANCO
    marco.line.color.rgb = RGBColor(0xD0, 0xD5, 0xDC)
    marco.line.width = Pt(0.75)
    marco.shadow.inherit = False
    marco.text_frame.paragraphs[0].text = ""

    try:
        from PIL import Image
        iw, ih = Image.open(io.BytesIO(imagen_bytes)).size
        escala = min(ancho / iw, alto_imagen / ih)
        w, h = int(iw * escala), int(ih * escala)
        x = left + (ancho - w) // 2
        y = top + (alto_imagen - h) // 2
        slide.shapes.add_picture(io.BytesIO(imagen_bytes), x, y, width=w, height=h)
    except Exception:
        slide.shapes.add_picture(io.BytesIO(imagen_bytes), left, top, width=ancho, height=alto_imagen)

    caja_pie = slide.shapes.add_textbox(left, top + alto_imagen, ancho, alto_pie)
    tf = caja_pie.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = _truncar(pie, 90) if pie else ""
    p.font.size, p.font.color.rgb = Pt(8), GRIS_TEXTO


def _slide_fotos(prs, datos, fotos=None, pagina=None, total_paginas=None):
    """Si `fotos` viene vacío/None: una slide con 6 casilleros en blanco
    (comportamiento original, para completar a mano). Si trae fotos (ya
    seleccionadas por el ajustador): esta función se llama una vez por
    página de hasta 6 fotos — ver generar_ficha_pptx."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitulo = datos.get("nickname") or datos.get("asegurado") or ""
    if total_paginas and total_paginas > 1:
        subtitulo = f"{subtitulo} — Página {pagina} de {total_paginas}"
    _agregar_header(slide, "REGISTRO FOTOGRÁFICO", subtitulo)

    posiciones = _posiciones_grilla_fotos()
    for i, (left, top, ancho, alto) in enumerate(posiciones, start=1):
        if fotos and i <= len(fotos):
            _casillero_con_foto(slide, left, top, ancho, alto, fotos[i - 1]["imagen"], fotos[i - 1].get("pie", ""))
        elif not fotos:
            _casillero_vacio(slide, left, top, ancho, alto, i)
    return slide


def _slide_gestiones(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _agregar_header(slide, "DETALLE DE GESTIONES REALIZADAS", datos.get("nickname") or datos.get("asegurado") or "")

    n_filas_dato = 6
    left, top = Inches(0.4), HEADER_ALTO + Emu(36576) + Inches(0.3)
    width = SLIDE_WIDTH - Inches(0.8)
    height = SLIDE_HEIGHT - top - Inches(0.3)

    graphic_frame = slide.shapes.add_table(n_filas_dato + 1, 2, left, top, width, height)
    tabla = graphic_frame.table
    tabla.columns[0].width = Inches(1.6)
    tabla.columns[1].width = width - Inches(1.6)

    for c, texto in enumerate(["Fecha", "Detalle de la gestión"]):
        celda = tabla.cell(0, c)
        celda.text_frame.paragraphs[0].text = texto
        celda.fill.solid()
        celda.fill.fore_color.rgb = NAVY_OSCURO
        run = celda.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(12), True, BLANCO

    for r in range(1, n_filas_dato + 1):
        for c in range(2):
            celda = tabla.cell(r, c)
            celda.fill.solid()
            celda.fill.fore_color.rgb = BLANCO
            celda.text_frame.paragraphs[0].text = ""
    return slide


def _slide_reserva(prs, datos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _agregar_header(slide, "RESERVA DEL CASO", datos.get("nickname") or datos.get("asegurado") or "")

    top = HEADER_ALTO + Emu(36576) + Inches(0.25)
    left, width = Inches(0.4), SLIDE_WIDTH - Inches(0.8)

    # --- Valor destacado (extraído de la planilla base) ---
    callout = slide.shapes.add_shape(1, left, top, width, Inches(0.85))
    callout.fill.solid()
    callout.fill.fore_color.rgb = NAVY
    callout.line.fill.background()
    callout.shadow.inherit = False
    tf = callout.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "RESERVA (Pérdida neta)"
    p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, TEAL
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = datos.get("perdida_neta_fmt") or "—"
    p2.font.size, p2.font.bold, p2.font.color.rgb = Pt(26), True, BLANCO

    # --- Desglose extraído de la planilla + columna de justificación manual ---
    top_tabla = top + Inches(1.05)
    height_tabla = SLIDE_HEIGHT - top_tabla - Inches(0.3)
    filas_valores = [
        ("Pérdida bruta", datos.get("perdida_bruta_fmt")),
        ("Deducible", datos.get("deducible_fmt")),
        ("Monto asegurado", datos.get("monto_asegurado_fmt")),
        ("Pérdida neta (Reserva)", datos.get("perdida_neta_fmt")),
        ("Gastos", datos.get("gastos_fmt")),
        ("Honorarios", datos.get("honorarios_fmt")),
    ]
    graphic_frame = slide.shapes.add_table(len(filas_valores) + 1, 3, left, top_tabla, width, height_tabla)
    tabla = graphic_frame.table
    tabla.columns[0].width = Inches(2.4)
    tabla.columns[1].width = Inches(1.8)
    tabla.columns[2].width = width - Inches(4.2)

    for c, texto in enumerate(["Concepto", "Monto (planilla)", "Justificación del ajustador"]):
        celda = tabla.cell(0, c)
        celda.text_frame.paragraphs[0].text = texto
        celda.fill.solid()
        celda.fill.fore_color.rgb = NAVY_OSCURO
        run = celda.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.color.rgb = Pt(11), True, BLANCO

    for r, (concepto, monto) in enumerate(filas_valores, start=1):
        celda_concepto = tabla.cell(r, 0)
        celda_concepto.text_frame.paragraphs[0].text = concepto
        celda_concepto.fill.solid()
        celda_concepto.fill.fore_color.rgb = GRIS_CLARO
        celda_concepto.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        celda_monto = tabla.cell(r, 1)
        celda_monto.text_frame.paragraphs[0].text = monto or "—"
        celda_monto.fill.solid()
        celda_monto.fill.fore_color.rgb = GRIS_CLARO
        celda_monto.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

        celda_just = tabla.cell(r, 2)
        celda_just.text_frame.paragraphs[0].text = ""
        celda_just.fill.solid()
        celda_just.fill.fore_color.rgb = BLANCO
    return slide


def _slide_texto_libre(prs, datos, titulo, texto_precargado=None, fuente=None):
    """Slide con un único cuadro grande de texto (usado por Estado Actual del
    Siniestro y Próximas Acciones). Si `texto_precargado` viene vacío, el
    cuadro queda en blanco con la nota "(Texto libre)"; si trae contenido
    (p.ej. la observación del Pipeline), se precarga ese texto —editable— y
    se indica su `fuente` en una nota al pie."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _agregar_header(slide, titulo, datos.get("nickname") or datos.get("asegurado") or "")

    top = HEADER_ALTO + Emu(36576) + Inches(0.3)
    left, width = Inches(0.4), SLIDE_WIDTH - Inches(0.8)
    alto_nota = Inches(0.3) if (texto_precargado and fuente) else Inches(0)
    height = SLIDE_HEIGHT - top - Inches(0.3) - alto_nota

    caja = slide.shapes.add_shape(1, left, top, width, height)
    caja.fill.solid()
    caja.fill.fore_color.rgb = GRIS_CLARO
    caja.line.color.rgb = RGBColor(0xB0, 0xB8, 0xC2)
    caja.line.width = Pt(1)
    caja.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    caja.shadow.inherit = False
    tf = caja.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left, tf.margin_right = Pt(10), Pt(10)
    tf.margin_top, tf.margin_bottom = Pt(8), Pt(8)
    p = tf.paragraphs[0]
    if texto_precargado:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p.text = texto_precargado
        p.font.size, p.font.color.rgb = Pt(13), GRIS_TEXTO
    else:
        p.text = "(Texto libre)"
        p.font.size, p.font.italic, p.font.color.rgb = Pt(11), True, RGBColor(0x8A, 0x93, 0x9E)

    if texto_precargado and fuente:
        nota = slide.shapes.add_textbox(left, top + height + Inches(0.03), width, alto_nota)
        pn = nota.text_frame.paragraphs[0]
        pn.text = f"Fuente: {fuente} — edítalo si corresponde."
        pn.font.size, pn.font.italic, pn.font.color.rgb = Pt(9), True, RGBColor(0x8A, 0x93, 0x9E)
    return slide


def _fmt_monto(valor, divisa):
    try:
        numero = float(valor)
    except (TypeError, ValueError):
        return None
    texto = f"{numero:,.0f}".replace(",", ".")
    return f"{texto} {divisa}".strip()


def generar_ficha_pptx(fila, observacion_pipeline=None, observacion_fuente=None, fotos_seleccionadas=None):
    """fila: dict o pandas.Series con los datos del caso, usando los mismos
    nombres de columna que trae la Base Maestra (ver app.py).

    observacion_pipeline: texto completo (opcional) para precargar la slide
    "Estado Actual del Siniestro" — normalmente la observación del Pipeline
    para ese Caso JPV. observacion_fuente describe de dónde salió (se
    muestra como nota al pie, p.ej. "Observaciones del Pipeline").

    fotos_seleccionadas: lista opcional de dicts {"imagen": bytes, "pie": str}
    (p.ej. de acta_inspeccion.extraer_fotos_acta, ya filtrada por el
    ajustador) para precargar el Registro Fotográfico. Se reparten en
    páginas de 6 fotos cada una; sin este argumento, queda una sola slide
    con 6 casilleros en blanco para completar a mano.
    """
    divisa = str(fila.get("Divisa") or "").strip()
    datos = dict(
        asegurado=fila.get("Asegurado"),
        nickname=fila.get("Nickname"),
        numero_siniestro=fila.get("Número de siniestro"),
        caso_jpv=fila.get("ID_Caso"),
        estado=fila.get("Estado_Actual") or fila.get("Estado"),
        fecha_ocurrencia=_fmt_fecha(fila.get("Fecha de ocurrencia")),
        fecha_denuncio=_fmt_fecha(fila.get("Fecha de denuncio")),
        fecha_asignacion=_fmt_fecha(fila.get("Fecha de asignación")),
        dias_asignacion=fila.get("Días desde asignación"),
        monto_asegurado_fmt=_fmt_monto(fila.get("Monto asegurado (en moneda del caso)"), divisa),
        perdida_bruta_fmt=_fmt_monto(fila.get("Perdida bruta (en moneda del caso)"), divisa),
        deducible_fmt=_fmt_monto(fila.get("Deducible (en moneda del caso)"), divisa),
        perdida_neta_fmt=_fmt_monto(fila.get("Perdida neta (en moneda del caso)"), divisa),
        gastos_fmt=_fmt_monto(fila.get("Gastos (UF)"), "UF"),
        honorarios_fmt=_fmt_monto(fila.get("Honorarios (UF)"), "UF"),
    )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _slide_resumen(prs, datos)
    _slide_descripcion(prs, datos)
    if fotos_seleccionadas:
        paginas = [
            fotos_seleccionadas[i:i + _FOTOS_POR_PAGINA]
            for i in range(0, len(fotos_seleccionadas), _FOTOS_POR_PAGINA)
        ]
        for num_pagina, pagina in enumerate(paginas, start=1):
            _slide_fotos(prs, datos, fotos=pagina, pagina=num_pagina, total_paginas=len(paginas))
    else:
        _slide_fotos(prs, datos)
    _slide_reserva(prs, datos)
    _slide_texto_libre(
        prs, datos, "ESTADO ACTUAL DEL SINIESTRO",
        texto_precargado=observacion_pipeline, fuente=observacion_fuente,
    )
    _slide_gestiones(prs, datos)
    _slide_texto_libre(prs, datos, "PRÓXIMAS ACCIONES")

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _fmt_fecha(valor):
    if valor in (None, "", "nan", "NaT"):
        return None
    fecha = pd.to_datetime(valor, errors="coerce")
    if pd.isna(fecha):
        return str(valor)
    return fecha.strftime("%d-%m-%Y")

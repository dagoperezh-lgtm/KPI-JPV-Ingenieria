"""
Generador del Reporte Ejecutivo "Estado de Cartera" para Marsh S.A.

Reutiliza el diseño de assets/marsh_template.pptx (7 slides) y reemplaza
únicamente los textos, tablas y el gráfico que contienen datos, preservando
el formato original (fuentes, colores, posiciones) del template.
"""
import io
import os
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "assets", "marsh_template.pptx")

MCL_UF = 5000
MCL_USD = 200000

MESES_ES = {
    1: "enero", 2: "febrero", 3: "marzo", 4: "abril", 5: "mayo", 6: "junio",
    7: "julio", 8: "agosto", 9: "septiembre", 10: "octubre", 11: "noviembre", 12: "diciembre",
}

# --- Mapa de shape_id por slide (0-indexado) construido a partir de assets/marsh_template.pptx ---
SLIDE1 = dict(fecha_corte=10, casos_activos=12, perdida_uf=15, perdida_usd=18)
SLIDE2 = dict(cartera_activa=11, division_split=17, divisa_split=23,
              exposicion_uf=29, exposicion_usd=35, mcl_pct=41, mcl_casos=43)
SLIDE3 = dict(mcl_count=11, mcl_pct_line=12, mcl_sum_line=13,
              otros_count=17, otros_pct_line=18, otros_sum_line=19,
              aseg1_nombre=21, aseg1_count=23, aseg2_nombre=24, aseg2_count=26,
              aseg3_nombre=27, aseg3_count=29, aseg4_nombre=30, aseg4_count=32,
              aseg_otros_count=35)
SLIDE4 = dict(tier100=12, tier75=18, tier50=24, tier0=30)
SLIDE5 = dict(caso_mas_antiguo=9, promedio_cartera=12, casos_600=15, tabla=17)
SLIDE6 = dict(intro=8, alerta=11, tabla=9)
SLIDE7 = dict(
    footer=32,
    pasos=[
        dict(titulo=10, desc=11),
        dict(titulo=15, desc=16),
        dict(titulo=20, desc=21),
        dict(titulo=30, desc=31),
        dict(titulo=38, desc=36),
    ],
)


def fmt_uf(valor):
    return f"{valor:,.0f}".replace(",", ".") + " UF"


def fmt_usd_m(valor):
    m = valor / 1_000_000
    s = f"{m:.1f}".replace(".", ",")
    return f"USD {s} M"


def fmt_fecha_larga(fecha):
    return f"{fecha.day} de {MESES_ES[fecha.month]} de {fecha.year}"


def filtrar_cartera_marsh(df_master, corredora_col="Corredora"):
    """Filtra los casos abiertos cuya Corredora contiene 'Marsh'."""
    if "Es_Abierto" not in df_master.columns:
        raise ValueError("df_master debe venir de procesar_datos_integrales()")
    df = df_master[df_master["Es_Abierto"]].copy()
    if corredora_col in df.columns:
        df = df[df[corredora_col].astype(str).str.contains("marsh", case=False, na=False)]
    return df


def preparar_tabla_casos(df_marsh_abiertos, fecha_corte):
    """Construye la tabla base (una fila por caso) con los campos que requiere
    el reporte. Los campos de juicio (Prob, Observación) quedan con valores
    por defecto para que el usuario los revise/edite antes de generar el pptx.
    """
    df = df_marsh_abiertos.copy()

    for col, default in [
        ("Divisa", "UF"),
        ("Perdida bruta (en moneda del caso)", 0),
        ("Asegurado", "Sin asegurado"),
        ("Nickname", ""),
        ("Última observación", ""),
    ]:
        if col not in df.columns:
            df[col] = default
        df[col] = df[col].fillna(default)

    divisa = df["Divisa"].astype(str).str.strip().str.upper()
    divisa = divisa.replace({"US$": "USD", "US $": "USD", "USD$": "USD", "U$": "USD"})

    perdida = pd.to_numeric(df["Perdida bruta (en moneda del caso)"], errors="coerce").fillna(0)
    dias = (pd.Timestamp(fecha_corte) - df["Fecha_Ingreso"]).dt.days.clip(lower=0).fillna(0).astype(int)
    mcl = ((divisa == "UF") & (perdida > MCL_UF)) | ((divisa == "USD") & (perdida > MCL_USD))

    nickname = df["Nickname"].astype(str).str.strip()
    asegurado = df["Asegurado"].astype(str).str.strip()
    nickname_default = nickname.where(nickname.ne("") & nickname.ne("None"), asegurado)

    tabla = pd.DataFrame({
        "Caso": df["ID_Caso"].astype(str),
        "Asegurado": asegurado,
        "Nickname": nickname_default,
        "Divisa": divisa,
        "Perdida_bruta": perdida,
        "Dias": dias,
        "Division": df["Area_Negocio"].astype(str),
        "MCL": mcl,
        "Observacion_sugerida": df["Última observación"].fillna("").astype(str),
        "Prob": 100,
        "Observacion": "",
    })
    return tabla.sort_values("Dias", ascending=False).reset_index(drop=True)


def calcular_kpis(tabla):
    """Calcula todos los KPIs objetivos del reporte a partir de la tabla
    (ya editada por el usuario, con Prob/Observacion definitivos)."""
    total = len(tabla)
    uf_mask = tabla["Divisa"] == "UF"
    usd_mask = tabla["Divisa"] == "USD"

    uf_count = int(uf_mask.sum())
    usd_count = int(usd_mask.sum())
    uf_total = float(tabla.loc[uf_mask, "Perdida_bruta"].sum())
    usd_total = float(tabla.loc[usd_mask, "Perdida_bruta"].sum())

    div_ing = int(tabla["Division"].str.contains("Ingenier", case=False, na=False).sum())
    div_otros = total - div_ing

    mcl = tabla[tabla["MCL"]]
    otros = tabla[~tabla["MCL"]]
    mcl_count = len(mcl)
    otros_count = len(otros)
    mcl_pct = round(100 * mcl_count / total) if total else 0
    otros_pct = 100 - mcl_pct if total else 0

    mcl_uf = float(mcl.loc[mcl["Divisa"] == "UF", "Perdida_bruta"].sum())
    mcl_usd = float(mcl.loc[mcl["Divisa"] == "USD", "Perdida_bruta"].sum())
    otros_uf = float(otros.loc[otros["Divisa"] == "UF", "Perdida_bruta"].sum())
    otros_usd = float(otros.loc[otros["Divisa"] == "USD", "Perdida_bruta"].sum())

    top_aseg = tabla.groupby("Asegurado").size().sort_values(ascending=False)
    top4 = list(top_aseg.head(4).items())
    while len(top4) < 4:
        top4.append(("—", 0))
    aseg_otros = int(top_aseg.iloc[4:].sum())

    dias_max = int(tabla["Dias"].max()) if total else 0
    dias_prom = round(tabla["Dias"].mean()) if total else 0
    dias_600 = int((tabla["Dias"] > 600).sum())

    prob = pd.to_numeric(tabla["Prob"], errors="coerce").fillna(100)
    tier_counts = {
        100: int((prob >= 90).sum()),
        75: int(((prob >= 60) & (prob < 90)).sum()),
        50: int(((prob >= 25) & (prob < 60)).sum()),
        0: int((prob < 25).sum()),
    }

    return dict(
        total=total, uf_count=uf_count, usd_count=usd_count,
        uf_total=uf_total, usd_total=usd_total,
        div_ing=div_ing, div_otros=div_otros,
        mcl_count=mcl_count, otros_count=otros_count,
        mcl_pct=mcl_pct, otros_pct=otros_pct,
        mcl_uf=mcl_uf, mcl_usd=mcl_usd, otros_uf=otros_uf, otros_usd=otros_usd,
        top4=top4, aseg_otros=aseg_otros,
        dias_max=dias_max, dias_prom=dias_prom, dias_600=dias_600,
        tier_counts=tier_counts,
    )


def _set_text(text_frame, texto):
    p = text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = texto
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = texto


def _set_shape_text(slide, shape_id, texto):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            _set_text(shape.text_frame, texto)
            return
    raise KeyError(f"shape_id {shape_id} no encontrado en el slide")


def _get_table(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.has_table:
            return shape.table
    raise KeyError(f"tabla shape_id {shape_id} no encontrada en el slide")


def _fill_table_rows(tabla_pptx, filas):
    """filas: lista de listas (cada una con el mismo número de columnas que la tabla).
    Rellena desde la fila 1 (la 0 es el encabezado) y limpia las filas sobrantes."""
    n_filas_datos = len(tabla_pptx.rows) - 1
    for i in range(n_filas_datos):
        for c in range(len(tabla_pptx.columns)):
            valor = filas[i][c] if i < len(filas) else ""
            _set_text(tabla_pptx.cell(i + 1, c).text_frame, str(valor))


def generar_pptx(fecha_corte, tabla, pasos, alerta_prioritaria):
    """
    fecha_corte: datetime.date/datetime
    tabla: DataFrame de preparar_tabla_casos() ya editado por el usuario
           (columnas Caso, Nickname, Divisa, Perdida_bruta, Dias, Division, MCL, Prob, Observacion)
    pasos: lista de hasta 5 dicts {'titulo':..., 'desc':...}
    alerta_prioritaria: texto libre para el callout del slide 6
    """
    kpis = calcular_kpis(tabla)
    prs = Presentation(TEMPLATE_PATH)

    # --- Slide 1: Portada ---
    s1 = prs.slides[0]
    _set_shape_text(s1, SLIDE1["fecha_corte"], fmt_fecha_larga(fecha_corte))
    _set_shape_text(s1, SLIDE1["casos_activos"], str(kpis["total"]))
    _set_shape_text(s1, SLIDE1["perdida_uf"], fmt_uf(kpis["uf_total"]))
    _set_shape_text(s1, SLIDE1["perdida_usd"], fmt_usd_m(kpis["usd_total"]))

    # --- Slide 2: Resumen general ---
    s2 = prs.slides[1]
    _set_shape_text(s2, SLIDE2["cartera_activa"], str(kpis["total"]))
    _set_shape_text(s2, SLIDE2["division_split"], f"{kpis['div_ing']} / {kpis['div_otros']}")
    _set_shape_text(s2, SLIDE2["divisa_split"], f"{kpis['uf_count']} UF · {kpis['usd_count']} USD")
    _set_shape_text(s2, SLIDE2["exposicion_uf"], fmt_uf(kpis["uf_total"]))
    _set_shape_text(s2, SLIDE2["exposicion_usd"], fmt_usd_m(kpis["usd_total"]))
    _set_shape_text(s2, SLIDE2["mcl_pct"], f"{kpis['mcl_pct']}% MCL")
    _set_shape_text(s2, SLIDE2["mcl_casos"], f"{kpis['mcl_count']} casos mayor cuantía")

    # --- Slide 3: Distribución MCL vs Otros ---
    s3 = prs.slides[2]
    _set_shape_text(s3, SLIDE3["mcl_count"], str(kpis["mcl_count"]))
    _set_shape_text(s3, SLIDE3["mcl_pct_line"], f"casos  ·  {kpis['mcl_pct']}% de la cartera")
    _set_shape_text(s3, SLIDE3["mcl_sum_line"], f"{fmt_uf(kpis['mcl_uf'])}  +  {fmt_usd_m(kpis['mcl_usd'])}")
    _set_shape_text(s3, SLIDE3["otros_count"], str(kpis["otros_count"]))
    _set_shape_text(s3, SLIDE3["otros_pct_line"], f"casos  ·  {kpis['otros_pct']}% de la cartera")
    _set_shape_text(s3, SLIDE3["otros_sum_line"], f"{fmt_uf(kpis['otros_uf'])}  +  {fmt_usd_m(kpis['otros_usd'])}")
    nombre_ids = ["aseg1_nombre", "aseg2_nombre", "aseg3_nombre", "aseg4_nombre"]
    count_ids = ["aseg1_count", "aseg2_count", "aseg3_count", "aseg4_count"]
    for (nombre, count), nombre_id, count_id in zip(kpis["top4"], nombre_ids, count_ids):
        _set_shape_text(s3, SLIDE3[nombre_id], str(nombre))
        _set_shape_text(s3, SLIDE3[count_id], str(count))
    _set_shape_text(s3, SLIDE3["aseg_otros_count"], str(kpis["aseg_otros"]))

    # --- Slide 4: Probabilidad de cierre ---
    s4 = prs.slides[3]
    tc = kpis["tier_counts"]
    _set_shape_text(s4, SLIDE4["tier100"], f"{tc[100]} caso" + ("s" if tc[100] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier75"], f"{tc[75]} caso" + ("s" if tc[75] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier50"], f"{tc[50]} caso" + ("s" if tc[50] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier0"], f"{tc[0]} caso" + ("s" if tc[0] != 1 else ""))
    for shape in s4.shapes:
        if shape.has_chart:
            cd = CategoryChartData()
            cd.categories = [
                f"100% – Cierta ({tc[100]})",
                f"75% – Alt. probable ({tc[75]})",
                f"50% – Podría ser ({tc[50]})",
                f"0% – Nula ({tc[0]})",
            ]
            cd.add_series("Prob. Cierre", (tc[100], tc[75], tc[50], tc[0]))
            shape.chart.replace_data(cd)

    # --- Slide 5: Tiempos de residencia (Top 10 más antiguos) ---
    s5 = prs.slides[4]
    _set_shape_text(s5, SLIDE5["caso_mas_antiguo"], f"{kpis['dias_max']} días")
    _set_shape_text(s5, SLIDE5["promedio_cartera"], f"~{kpis['dias_prom']} días")
    _set_shape_text(s5, SLIDE5["casos_600"], f"{kpis['dias_600']} caso" + ("s" if kpis["dias_600"] != 1 else ""))
    top10_antiguos = tabla.sort_values("Dias", ascending=False).head(10)
    filas = []
    for i, (_, row) in enumerate(top10_antiguos.iterrows(), start=1):
        filas.append([
            i, row["Caso"], row["Nickname"], row["Divisa"], row["Dias"],
            f"{int(row['Prob'])}%", row["Observacion"] or row["Observacion_sugerida"],
        ])
    _fill_table_rows(_get_table(s5, SLIDE5["tabla"]), filas)

    # --- Slide 6: Casos que requieren atención especial (Prob < 100%) ---
    s6 = prs.slides[5]
    atencion = tabla[pd.to_numeric(tabla["Prob"], errors="coerce").fillna(100) < 100]
    atencion = atencion.sort_values("Perdida_bruta", ascending=False).head(10)
    _set_shape_text(
        s6, SLIDE6["intro"],
        f"{len(atencion)} casos con probabilidad de cierre inferior al 100% acumulan exposición "
        f"significativa. Ordenados por magnitud de pérdida bruta.",
    )
    _set_shape_text(s6, SLIDE6["alerta"], alerta_prioritaria)
    filas = []
    for _, row in atencion.iterrows():
        filas.append([
            row["Caso"], row["Nickname"], f"{int(row['Prob'])}%", row["Divisa"],
            "✓" if row["MCL"] else "", row["Observacion"] or row["Observacion_sugerida"],
        ])
    _fill_table_rows(_get_table(s6, SLIDE6["tabla"]), filas)

    # --- Slide 7: Próximos pasos y focos de gestión ---
    s7 = prs.slides[6]
    for i, ids in enumerate(SLIDE7["pasos"]):
        paso = pasos[i] if i < len(pasos) else {"titulo": "", "desc": ""}
        _set_shape_text(s7, ids["titulo"], paso.get("titulo", ""))
        _set_shape_text(s7, ids["desc"], paso.get("desc", ""))
    _set_shape_text(
        s7, SLIDE7["footer"],
        f"JPV Asociados · Ajustadores Especializados  ·  Análisis al {fmt_fecha_larga(fecha_corte)}  ·  Cartera MARSH S.A.",
    )

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

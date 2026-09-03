"""
Generador del Reporte Ejecutivo "Estado de Cartera".

Reutiliza el diseño de assets/plantilla_estado_cartera.pptx (7 slides) y
reemplaza únicamente los textos, tablas y el gráfico que contienen datos,
preservando el formato original (fuentes, colores, posiciones) del template.
Sirve para cualquier corte de cartera: por Corredora, Compañía de seguros
(Aseguradora), Asegurado, o una combinación de estos filtros.
"""
import copy
import io
import os
import re
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

try:
    import gspread
    from google.oauth2.service_account import Credentials
except ImportError:
    gspread = None
    Credentials = None

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "assets", "plantilla_estado_cartera.pptx")

# Pipeline actualizado por el equipo en Google Sheets, compartido en modo
# Lector con la cuenta de servicio opscontrol-bot@... (misma cuenta que ya
# usa OpsControl para la Base Maestra en tablero_datos.py). Los datos viven
# en la pestaña Pipeline_Backup, con la misma estructura que Base_Maestra:
# fila 1 = metadata (FECHA_ACTUALIZACION), fila 2 = encabezados reales.
PIPELINE_SHEET_URL = "https://docs.google.com/spreadsheets/d/1aZYWn1JI_w4S9mQFpaCJRAQfr6E0iKqXMcMJ_xxlaNI/edit?usp=sharing"
PIPELINE_WORKSHEET_NAME = "Pipeline_Backup"
PIPELINE_HEADER_ROW = 2

MCL_UF = 5000
MCL_USD = 200000

# La columna "Indicación Probabilidad" del Pipeline usa estas 4 categorías.
INDICACION_A_PROB = {
    "cierta": 100,
    "altamente probable": 75,
    "podría ser": 50,
    "podria ser": 50,
    "nula": 0,
}

_TIMESTAMP_PREFIJO_RE = re.compile(r"^\s*\[\d{1,2}/\d{1,2}/\d{2,4}[^\]]*\]\s*")

MESES_ES = {
    1: "enero", 2: "febrero", 3: "marzo", 4: "abril", 5: "mayo", 6: "junio",
    7: "julio", 8: "agosto", 9: "septiembre", 10: "octubre", 11: "noviembre", 12: "diciembre",
}

# --- Mapa de shape_id por slide (0-indexado) construido a partir de assets/plantilla_estado_cartera.pptx ---
SLIDE1 = dict(titulo_cartera=6, fecha_corte=10, casos_activos=12, perdida_uf=15, perdida_usd=18)
SLIDE2 = dict(titulo_cartera=6, cartera_activa=11, division_split=17, divisa_split=23,
              exposicion_uf=29, exposicion_usd=35, mcl_pct=41, mcl_casos=43)
SLIDE3 = dict(mcl_count=11, mcl_pct_line=12, mcl_sum_line=13,
              otros_count=17, otros_pct_line=18, otros_sum_line=19,
              aseg1_nombre=21, aseg1_barra=22, aseg1_count=23,
              aseg2_nombre=24, aseg2_barra=25, aseg2_count=26,
              aseg3_nombre=27, aseg3_barra=28, aseg3_count=29,
              aseg4_nombre=30, aseg4_barra=31, aseg4_count=32,
              aseg_otros_barra=34, aseg_otros_count=35)
SLIDE4 = dict(tier100=12, tier75=18, tier50=24, tier0=30)
SLIDE5 = dict(caso_mas_antiguo=9, promedio_cartera=12, promedio_cartera_label=13, casos_600=15, tabla=17)
SLIDE6 = dict(intro=8, alerta=11, tabla=9)
# Slide "Gestiones Iniciales" (penúltima): clon exacto de la slide 5, mismos
# shape_id, pero con los 3 KPI y la tabla reutilizados para casos recientes.
SLIDE_GESTIONES = dict(
    titulo=5, subtitulo=6,
    kpi1_valor=9, kpi1_label=10,
    kpi2_valor=12, kpi2_label=13,
    kpi3_valor=15, kpi3_label=16,
    tabla=17,
)
SLIDE7 = dict(
    footer=32,
    pasos=[
        dict(titulo=10, desc=11),
        dict(titulo=15, desc=16),
        dict(titulo=20, desc=21),
        dict(titulo=30, desc=31),
        dict(titulo=38, desc=36),
    ],
)


def fmt_uf(valor):
    return f"{valor:,.0f}".replace(",", ".") + " UF"


def fmt_usd_m(valor):
    m = valor / 1_000_000
    s = f"{m:.1f}".replace(".", ",")
    return f"USD {s} M"


def fmt_fecha_larga(fecha):
    return f"{fecha.day} de {MESES_ES[fecha.month]} de {fecha.year}"


def _fila_siniestro_nickname(siniestro, nickname, max_chars=70):
    """Combina el N° de siniestro con el nickname para las columnas 'Siniestro
    / Nickname' de las tablas (slides 5, 6 y 'Gestiones Iniciales'). El N° de
    siniestro nunca se trunca (es un identificador); el nickname se recorta
    con el espacio restante para no salirse del ancho de la columna."""
    siniestro = str(siniestro).strip()
    nickname = str(nickname).strip()
    if not siniestro or siniestro.lower() in ("nan", "none"):
        return _truncar_texto(nickname, max_chars)
    prefijo = f"Sin. {siniestro} — "
    espacio_nickname = max(15, max_chars - len(prefijo))
    return prefijo + _truncar_texto(nickname, espacio_nickname)


def _truncar_texto(texto, max_chars):
    """Recorta observaciones largas para que la fila no crezca más de lo que
    permite la plantilla (si no, las últimas filas de la tabla quedan fuera
    del área visible de la slide en modo presentación)."""
    texto = str(texto).strip()
    if len(texto) <= max_chars:
        return texto
    return texto[: max_chars - 1].rstrip() + "…"


def cargar_pipeline(archivo, sheet_name=None):
    """Lee el Excel de Pipeline (una fila por caso activo, con Probabilidad de
    cierre y Observaciones ya cargadas por el equipo)."""
    xl = pd.ExcelFile(archivo)
    hoja = sheet_name or xl.sheet_names[0]
    return pd.read_excel(archivo, sheet_name=hoja)


def _cliente_google_sheets(credenciales_sa):
    if gspread is None:
        raise RuntimeError("Faltan las librerías 'gspread' y 'google-auth' (revisa requirements.txt).")
    scope = [
        "https://www.googleapis.com/auth/spreadsheets.readonly",
        "https://www.googleapis.com/auth/drive.readonly",
    ]
    creds = Credentials.from_service_account_info(credenciales_sa, scopes=scope)
    return gspread.authorize(creds)


def cargar_hoja_desde_google_sheets(credenciales_sa, worksheet_name, header_row=1, sheet_url=None):
    """Descarga cualquier pestaña del Google Sheet de OpsControl (Pipeline_Backup,
    Base_Maestra, etc.) usando una cuenta de servicio (st.secrets["gcp_service_account"]).
    No requiere que la hoja esté compartida públicamente: basta con compartirla
    como Lector con el email de la cuenta de servicio.

    credenciales_sa: dict con el JSON de la cuenta de servicio.
    worksheet_name: nombre de la pestaña a leer.
    header_row: fila donde están los encabezados reales (1 si es la primera
    fila; 2 si, como Base_Maestra o Pipeline_Backup, la fila 1 es metadata).
    sheet_url: enlace del Google Sheet (por defecto PIPELINE_SHEET_URL).
    """
    client = _cliente_google_sheets(credenciales_sa)
    doc = client.open_by_url(sheet_url or PIPELINE_SHEET_URL)
    hoja = doc.worksheet(worksheet_name)
    return pd.DataFrame(hoja.get_all_records(head=header_row))


def cargar_pipeline_desde_google_sheets(credenciales_sa, sheet_url=None, worksheet_name=None):
    """Descarga el Pipeline directamente desde Google Sheets (ver
    cargar_hoja_desde_google_sheets)."""
    return cargar_hoja_desde_google_sheets(
        credenciales_sa,
        worksheet_name or PIPELINE_WORKSHEET_NAME,
        header_row=PIPELINE_HEADER_ROW,
        sheet_url=sheet_url,
    )


def filtrar_pipeline(df_pipeline, corredoras=None, aseguradoras=None, asegurados=None, ajustadores=None):
    """Filtra el Pipeline por Corredora / Compañía de seguros / Asegurado / Ajustador senior.

    Cada parámetro es una lista opcional de valores exactos (tal como aparecen
    en el archivo). Si una lista viene vacía o None, ese filtro no se aplica.
    Cuando se combinan varios filtros, se aplican con AND (intersección).
    """
    df = df_pipeline.copy()
    if corredoras and "Corredora" in df.columns:
        df = df[df["Corredora"].isin(corredoras)]
    if aseguradoras and "Compañía de seguros" in df.columns:
        df = df[df["Compañía de seguros"].isin(aseguradoras)]
    if asegurados and "Asegurado" in df.columns:
        df = df[df["Asegurado"].isin(asegurados)]
    if ajustadores and "Ajustador senior" in df.columns:
        df = df[df["Ajustador senior"].isin(ajustadores)]
    return df


def sugerir_titulo_cartera(corredoras=None, aseguradoras=None, asegurados=None):
    """Sugiere un título de portada a partir de los filtros activos (editable por el usuario)."""
    if asegurados:
        return " · ".join(asegurados)
    partes = []
    if corredoras:
        partes.append(" · ".join(corredoras))
    if aseguradoras:
        partes.append(" · ".join(aseguradoras))
    return " · ".join(partes) if partes else "Cartera General"


def sugerir_titulo_ajustador(ajustadores=None):
    """Sugiere un título de portada a partir de los Ajustadores senior seleccionados."""
    return " · ".join(ajustadores) if ajustadores else "Cartera General"


def limpiar_observacion(texto):
    return _TIMESTAMP_PREFIJO_RE.sub("", str(texto)).strip()


def preparar_tabla_casos(df_pipeline_filtrado, fecha_corte):
    """Construye la tabla base (una fila por caso) a partir del Pipeline.
    Probabilidad de cierre y Observación ya vienen cargadas por el equipo
    (columnas 'Indicación Probabilidad' / 'Probabilidad cierre 2026' y
    'Observaciones'); el usuario solo las revisa/ajusta antes de generar el pptx.
    """
    df = df_pipeline_filtrado.copy()

    for col, default in [
        ("Divisa", "UF"),
        ("Perdida bruta (en moneda del caso)", 0),
        ("Asegurado", "Sin asegurado"),
        ("Nickname", ""),
        ("Número de siniestro", ""),
        ("Observaciones", ""),
        ("Contenido último movimiento", ""),
        ("División", "Sin División"),
    ]:
        if col not in df.columns:
            df[col] = default
        df[col] = df[col].fillna(default)

    divisa = df["Divisa"].astype(str).str.strip().str.upper()
    divisa = divisa.replace({"US$": "USD", "US $": "USD", "USD$": "USD", "U$": "USD"})

    perdida = pd.to_numeric(df["Perdida bruta (en moneda del caso)"], errors="coerce").fillna(0)
    fecha_ingreso = pd.to_datetime(df["Creado en"], errors="coerce")
    dias = (pd.Timestamp(fecha_corte) - fecha_ingreso).dt.days.clip(lower=0).fillna(0).astype(int)
    mcl = ((divisa == "UF") & (perdida > MCL_UF)) | ((divisa == "USD") & (perdida > MCL_USD))

    nickname = df["Nickname"].astype(str).str.strip()
    asegurado = df["Asegurado"].astype(str).str.strip()
    nickname_default = nickname.where(nickname.ne("") & nickname.ne("None") & nickname.ne("nan"), asegurado)

    # Puede llegar como número (float, con ".0" al final) si Excel/Sheets lo
    # detectó como columna numérica; se normaliza siempre a texto limpio.
    siniestro = df["Número de siniestro"].astype(str).str.strip()
    siniestro = siniestro.str.replace(r"\.0$", "", regex=True)
    siniestro = siniestro.where(~siniestro.str.lower().isin(["nan", "none"]), "")

    indicacion = df.get("Indicación Probabilidad", pd.Series("", index=df.index)).astype(str).str.strip().str.lower()
    prob_indicacion = indicacion.map(INDICACION_A_PROB)
    prob_numerica = pd.to_numeric(df.get("Probabilidad cierre 2026", pd.Series(1.0, index=df.index)), errors="coerce") * 100
    prob = prob_indicacion.fillna(prob_numerica).fillna(100).round().astype(int)

    observacion = df["Observaciones"].astype(str).apply(limpiar_observacion)
    fallback = df["Contenido último movimiento"].astype(str).apply(limpiar_observacion)
    observacion_final = observacion.where(observacion.ne("") & observacion.ne("nan"), fallback)

    tabla = pd.DataFrame({
        "Caso": df["Número de caso"].astype(str),
        "Numero_Siniestro": siniestro,
        "Asegurado": asegurado,
        "Nickname": nickname_default,
        "Divisa": divisa,
        "Perdida_bruta": perdida,
        "Dias": dias,
        "Division": df["División"].astype(str),
        "MCL": mcl,
        "Observacion_sugerida": observacion_final,
        "Prob": prob,
        "Observacion": observacion_final,
    })
    return tabla.sort_values("Dias", ascending=False).reset_index(drop=True)


def calcular_kpis(tabla):
    """Calcula todos los KPIs objetivos del reporte a partir de la tabla
    (ya editada por el usuario, con Prob/Observacion definitivos)."""
    total = len(tabla)
    uf_mask = tabla["Divisa"] == "UF"
    usd_mask = tabla["Divisa"] == "USD"

    uf_count = int(uf_mask.sum())
    usd_count = int(usd_mask.sum())
    uf_total = float(tabla.loc[uf_mask, "Perdida_bruta"].sum())
    usd_total = float(tabla.loc[usd_mask, "Perdida_bruta"].sum())

    div_ing = int(tabla["Division"].str.contains("Ingenier", case=False, na=False).sum())
    div_otros = total - div_ing

    mcl = tabla[tabla["MCL"]]
    otros = tabla[~tabla["MCL"]]
    mcl_count = len(mcl)
    otros_count = len(otros)
    mcl_pct = round(100 * mcl_count / total) if total else 0
    otros_pct = 100 - mcl_pct if total else 0

    mcl_uf = float(mcl.loc[mcl["Divisa"] == "UF", "Perdida_bruta"].sum())
    mcl_usd = float(mcl.loc[mcl["Divisa"] == "USD", "Perdida_bruta"].sum())
    otros_uf = float(otros.loc[otros["Divisa"] == "UF", "Perdida_bruta"].sum())
    otros_usd = float(otros.loc[otros["Divisa"] == "USD", "Perdida_bruta"].sum())

    top_aseg = tabla.groupby("Asegurado").size().sort_values(ascending=False)
    top4 = list(top_aseg.head(4).items())
    while len(top4) < 4:
        top4.append(("—", 0))
    aseg_otros = int(top_aseg.iloc[4:].sum())

    dias_max = int(tabla["Dias"].max()) if total else 0
    dias_prom = round(tabla["Dias"].mean()) if total else 0
    dias_600 = int((tabla["Dias"] > 600).sum())

    prob = pd.to_numeric(tabla["Prob"], errors="coerce").fillna(100)
    tier_counts = {
        100: int((prob >= 90).sum()),
        75: int(((prob >= 60) & (prob < 90)).sum()),
        50: int(((prob >= 25) & (prob < 60)).sum()),
        0: int((prob < 25).sum()),
    }

    return dict(
        total=total, uf_count=uf_count, usd_count=usd_count,
        uf_total=uf_total, usd_total=usd_total,
        div_ing=div_ing, div_otros=div_otros,
        mcl_count=mcl_count, otros_count=otros_count,
        mcl_pct=mcl_pct, otros_pct=otros_pct,
        mcl_uf=mcl_uf, mcl_usd=mcl_usd, otros_uf=otros_uf, otros_usd=otros_usd,
        top4=top4, aseg_otros=aseg_otros,
        dias_max=dias_max, dias_prom=dias_prom, dias_600=dias_600,
        tier_counts=tier_counts,
    )


def _set_text(text_frame, texto):
    p = text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = texto
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = texto


def _get_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise KeyError(f"shape_id {shape_id} no encontrado en el slide")


def _set_shape_text(slide, shape_id, texto):
    _set_text(_get_shape(slide, shape_id).text_frame, texto)


def _get_table(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id and shape.has_table:
            return shape.table
    raise KeyError(f"tabla shape_id {shape_id} no encontrada en el slide")


def _fill_table_rows(tabla_pptx, filas):
    """filas: lista de listas (cada una con el mismo número de columnas que la tabla).
    Rellena desde la fila 1 (la 0 es el encabezado) y limpia las filas sobrantes."""
    n_filas_datos = len(tabla_pptx.rows) - 1
    for i in range(n_filas_datos):
        for c in range(len(tabla_pptx.columns)):
            valor = filas[i][c] if i < len(filas) else ""
            _set_text(tabla_pptx.cell(i + 1, c).text_frame, str(valor))


# Mismo cuadro que ocupaba la tabla original de 10 filas en la slide de
# Tiempos de Residencia (assets/plantilla_estado_cartera.pptx, shape_id 17).
_TABLA_DETALLE_LEFT = 256032
_TABLA_DETALLE_TOP = 1622909
_TABLA_DETALLE_WIDTH = 7644384
_TABLA_DETALLE_HEIGHT = 3323995
_TABLA_DETALLE_FILAS_MAX = 20
_TABLA_DETALLE_COLS = [
    # (encabezado, ancho, max_chars para truncar el valor)
    ("#", Inches(0.35), 3),
    ("Caso", Inches(0.75), 12),
    ("Siniestro", Inches(1.0), 18),
    ("Nickname / Asegurado", Inches(4.76), 95),
    ("Divisa", Inches(0.5), 6),
    ("Días", Inches(0.5), 6),
    ("Prob.", Inches(0.5), 6),
]


def _quitar_shape(slide, shape_id):
    shape = _get_shape(slide, shape_id)
    shape._element.getparent().remove(shape._element)


def _construir_tabla_detalle(slide, filas):
    """Tabla más densa que las demás del reporte (hasta 20 filas de datos en
    el mismo espacio donde las otras tablas ponen 10), pensada solo para
    listar la cartera completa. Se construye desde cero (no es un clon de
    ninguna tabla existente) para poder controlar el alto de fila y el
    tamaño de fuente sin afectar al resto de las slides."""
    n_filas = _TABLA_DETALLE_FILAS_MAX + 1
    n_cols = len(_TABLA_DETALLE_COLS)
    graphic_frame = slide.shapes.add_table(
        n_filas, n_cols, _TABLA_DETALLE_LEFT, _TABLA_DETALLE_TOP, _TABLA_DETALLE_WIDTH, _TABLA_DETALLE_HEIGHT
    )
    tabla_pptx = graphic_frame.table
    tabla_pptx.horz_banding = False

    for col, (_, ancho, _) in zip(tabla_pptx.columns, _TABLA_DETALLE_COLS):
        col.width = ancho

    alto_header = Inches(0.28)
    alto_dato = Inches((_TABLA_DETALLE_HEIGHT / 914400 - 0.28) / _TABLA_DETALLE_FILAS_MAX)
    tabla_pptx.rows[0].height = alto_header
    for fila in list(tabla_pptx.rows)[1:]:
        fila.height = alto_dato

    for c, (encabezado, _, _) in enumerate(_TABLA_DETALLE_COLS):
        cell = tabla_pptx.cell(0, c)
        _set_text(cell.text_frame, encabezado)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0D, 0x1F, 0x38)
        cell.margin_left, cell.margin_right = Pt(3), Pt(3)
        cell.margin_top, cell.margin_bottom = Pt(1), Pt(1)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold = Pt(8), True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i in range(_TABLA_DETALLE_FILAS_MAX):
        for c, (_, _, max_chars) in enumerate(_TABLA_DETALLE_COLS):
            valor = _truncar_texto(filas[i][c], max_chars) if i < len(filas) else ""
            cell = tabla_pptx.cell(i + 1, c)
            cell.fill.background()
            cell.margin_left, cell.margin_right = Pt(3), Pt(3)
            cell.margin_top, cell.margin_bottom = Pt(0.5), Pt(0.5)
            cell.text_frame.word_wrap = False
            _set_text(cell.text_frame, valor)
            if cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(6.5)
    return tabla_pptx


def _duplicar_slide(prs, index):
    """Clona la slide en `index` (mismo layout, shapes, tabla e imagen) y la
    agrega al final de la presentación. Los shape_id quedan idénticos a los
    de la slide original, así que se puede rellenar con el mismo mapa de
    shape_id (p.ej. SLIDE_GESTIONES) sin importar cuántas copias se hagan."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    rid_map = {}
    for shape in source.shapes:
        el = copy.deepcopy(shape._element)
        for blip in el.findall(".//" + qn("a:blip")):
            old_rid = blip.get(qn("r:embed"))
            if old_rid:
                if old_rid not in rid_map:
                    image_part = source.part.rels[old_rid].target_part
                    rid_map[old_rid] = dest.part.relate_to(image_part, source.part.rels[old_rid].reltype)
                blip.set(qn("r:embed"), rid_map[old_rid])
        dest.shapes._spTree.append(el)
    return dest


def _indice_por_slide_id(prs, slide_id):
    for i, s in enumerate(prs.slides):
        if s.slide_id == slide_id:
            return i
    raise ValueError(f"slide_id {slide_id} no encontrado")


def _duplicar_slide_antes_de(prs, index_fuente, slide_id_referencia):
    """Clona la slide `index_fuente` y la inserta justo antes de la slide con
    `slide_id_referencia` (recalculado en cada llamada, así que insertar
    varias veces seguidas las deja en el orden correcto, una tras otra)."""
    nueva = _duplicar_slide(prs, index_fuente)
    xml_slides = prs.slides._sldIdLst
    elemento = list(xml_slides)[-1]  # la recién agregada, siempre queda última
    xml_slides.remove(elemento)
    xml_slides.insert(_indice_por_slide_id(prs, slide_id_referencia), elemento)
    return nueva


def generar_pptx(fecha_corte, titulo_cartera, tabla, pasos, alerta_prioritaria):
    """
    fecha_corte: datetime.date/datetime
    titulo_cartera: texto para la portada/resumen/pie (ej. 'MARSH S.A.', 'Essbio S.A.',
                    'Aon Risk Services · Essbio S.A.', 'Cartera General', etc.)
    tabla: DataFrame de preparar_tabla_casos() ya editado por el usuario
           (columnas Caso, Nickname, Divisa, Perdida_bruta, Dias, Division, MCL, Prob, Observacion)
    pasos: lista de hasta 5 dicts {'titulo':..., 'desc':...}
    alerta_prioritaria: texto libre para el callout del slide 6
    """
    kpis = calcular_kpis(tabla)
    prs = Presentation(TEMPLATE_PATH)

    # --- Slide 1: Portada ---
    s1 = prs.slides[0]
    _set_shape_text(s1, SLIDE1["titulo_cartera"], titulo_cartera)
    _set_shape_text(s1, SLIDE1["fecha_corte"], fmt_fecha_larga(fecha_corte))
    _set_shape_text(s1, SLIDE1["casos_activos"], str(kpis["total"]))
    _set_shape_text(s1, SLIDE1["perdida_uf"], fmt_uf(kpis["uf_total"]))
    _set_shape_text(s1, SLIDE1["perdida_usd"], fmt_usd_m(kpis["usd_total"]))

    # --- Slide 2: Resumen general ---
    s2 = prs.slides[1]
    _set_shape_text(s2, SLIDE2["titulo_cartera"], titulo_cartera)
    _set_shape_text(s2, SLIDE2["cartera_activa"], str(kpis["total"]))
    _set_shape_text(s2, SLIDE2["division_split"], f"{kpis['div_ing']} / {kpis['div_otros']}")
    _set_shape_text(s2, SLIDE2["divisa_split"], f"{kpis['uf_count']} UF · {kpis['usd_count']} USD")
    _set_shape_text(s2, SLIDE2["exposicion_uf"], fmt_uf(kpis["uf_total"]))
    _set_shape_text(s2, SLIDE2["exposicion_usd"], fmt_usd_m(kpis["usd_total"]))
    _set_shape_text(s2, SLIDE2["mcl_pct"], f"{kpis['mcl_pct']}% MCL")
    _set_shape_text(s2, SLIDE2["mcl_casos"], f"{kpis['mcl_count']} casos mayor cuantía")

    # --- Slide 3: Distribución MCL vs Otros ---
    s3 = prs.slides[2]
    _set_shape_text(s3, SLIDE3["mcl_count"], str(kpis["mcl_count"]))
    _set_shape_text(s3, SLIDE3["mcl_pct_line"], f"casos  ·  {kpis['mcl_pct']}% de la cartera")
    _set_shape_text(s3, SLIDE3["mcl_sum_line"], f"{fmt_uf(kpis['mcl_uf'])}  +  {fmt_usd_m(kpis['mcl_usd'])}")
    _set_shape_text(s3, SLIDE3["otros_count"], str(kpis["otros_count"]))
    _set_shape_text(s3, SLIDE3["otros_pct_line"], f"casos  ·  {kpis['otros_pct']}% de la cartera")
    _set_shape_text(s3, SLIDE3["otros_sum_line"], f"{fmt_uf(kpis['otros_uf'])}  +  {fmt_usd_m(kpis['otros_usd'])}")
    nombre_ids = ["aseg1_nombre", "aseg2_nombre", "aseg3_nombre", "aseg4_nombre"]
    barra_ids = ["aseg1_barra", "aseg2_barra", "aseg3_barra", "aseg4_barra"]
    count_ids = ["aseg1_count", "aseg2_count", "aseg3_count", "aseg4_count"]

    # El orden (Top 4 + Otros) y el largo de la barra los define la misma
    # cantidad de casos que muestra el número — evita sugerir un monto de
    # pérdida/pago sobre casos cuya liquidación todavía no termina.
    BARRA_IZQUIERDA = 2926080
    BARRA_ANCHO_MAX = 5120640
    BARRA_ANCHO_MIN = 91440
    BARRA_ESPACIO_TEXTO = 91440
    max_count = max([c for _, c in kpis["top4"]] + [kpis["aseg_otros"]], default=0)

    def _ancho_barra(valor_count):
        if max_count <= 0:
            return BARRA_ANCHO_MIN
        return max(BARRA_ANCHO_MIN, round(valor_count / max_count * BARRA_ANCHO_MAX))

    def _dibujar_barra(barra_id, count_id, valor_count):
        ancho = _ancho_barra(valor_count)
        _get_shape(s3, SLIDE3[barra_id]).width = ancho
        shape_count = _get_shape(s3, SLIDE3[count_id])
        shape_count.left = BARRA_IZQUIERDA + ancho + BARRA_ESPACIO_TEXTO
        _set_text(shape_count.text_frame, str(valor_count))

    for (nombre, count), nombre_id, barra_id, count_id in zip(kpis["top4"], nombre_ids, barra_ids, count_ids):
        _set_shape_text(s3, SLIDE3[nombre_id], str(nombre))
        _dibujar_barra(barra_id, count_id, count)
    _dibujar_barra("aseg_otros_barra", "aseg_otros_count", kpis["aseg_otros"])

    # --- Slide 4: Probabilidad de cierre ---
    s4 = prs.slides[3]
    tc = kpis["tier_counts"]
    _set_shape_text(s4, SLIDE4["tier100"], f"{tc[100]} caso" + ("s" if tc[100] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier75"], f"{tc[75]} caso" + ("s" if tc[75] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier50"], f"{tc[50]} caso" + ("s" if tc[50] != 1 else ""))
    _set_shape_text(s4, SLIDE4["tier0"], f"{tc[0]} caso" + ("s" if tc[0] != 1 else ""))
    for shape in s4.shapes:
        if shape.has_chart:
            cd = CategoryChartData()
            cd.categories = [
                f"100% – Cierta ({tc[100]})",
                f"75% – Alt. probable ({tc[75]})",
                f"50% – Podría ser ({tc[50]})",
                f"0% – Nula ({tc[0]})",
            ]
            cd.add_series("Prob. Cierre", (tc[100], tc[75], tc[50], tc[0]))
            shape.chart.replace_data(cd)

    # --- Slide 5: Tiempos de residencia (Top 10 más antiguos) ---
    s5 = prs.slides[4]
    _set_shape_text(s5, SLIDE5["caso_mas_antiguo"], f"{kpis['dias_max']} días")
    _set_shape_text(s5, SLIDE5["promedio_cartera"], f"~{kpis['dias_prom']} días")
    etiqueta_promedio = "Promedio cartera" if titulo_cartera == "Cartera General" else f"Promedio cartera {titulo_cartera}"
    _set_shape_text(s5, SLIDE5["promedio_cartera_label"], etiqueta_promedio)
    _set_shape_text(s5, SLIDE5["casos_600"], f"{kpis['dias_600']} caso" + ("s" if kpis["dias_600"] != 1 else ""))
    top10_antiguos = tabla.sort_values("Dias", ascending=False).head(10)
    filas = []
    for i, (_, row) in enumerate(top10_antiguos.iterrows(), start=1):
        observacion = row["Observacion"] or row["Observacion_sugerida"]
        filas.append([
            i, row["Caso"], _fila_siniestro_nickname(row["Numero_Siniestro"], row["Nickname"]), row["Divisa"], row["Dias"],
            f"{int(row['Prob'])}%", _truncar_texto(observacion, 100),
        ])
    _fill_table_rows(_get_table(s5, SLIDE5["tabla"]), filas)

    # --- Slide 6: Casos que requieren atención especial (Prob < 100%) ---
    s6 = prs.slides[5]
    atencion = tabla[pd.to_numeric(tabla["Prob"], errors="coerce").fillna(100) < 100]
    atencion = atencion.sort_values("Perdida_bruta", ascending=False).head(10)
    _set_shape_text(
        s6, SLIDE6["intro"],
        f"{len(atencion)} casos con probabilidad de cierre inferior al 100% acumulan exposición "
        f"significativa. Ordenados por magnitud de pérdida bruta.",
    )
    _set_shape_text(s6, SLIDE6["alerta"], alerta_prioritaria)
    filas = []
    for _, row in atencion.iterrows():
        observacion = row["Observacion"] or row["Observacion_sugerida"]
        filas.append([
            row["Caso"], _fila_siniestro_nickname(row["Numero_Siniestro"], row["Nickname"]), f"{int(row['Prob'])}%", row["Divisa"],
            "✓" if row["MCL"] else "", _truncar_texto(observacion, 130),
        ])
    _fill_table_rows(_get_table(s6, SLIDE6["tabla"]), filas)

    # --- Slide "Gestiones Iniciales" (penúltima): casos ingresados hace menos de 2 semanas ---
    s_gestiones = prs.slides[6]
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["titulo"], "GESTIONES INICIALES")
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["subtitulo"], "Casos ingresados hace menos de 2 semanas")

    recientes = tabla[tabla["Dias"] < 14].sort_values("Dias", ascending=False)
    sin_gestion = recientes[
        (recientes["Observacion"].astype(str).str.strip() == "")
        & (recientes["Observacion_sugerida"].astype(str).str.strip() == "")
    ]
    dias_prom_recientes = round(recientes["Dias"].mean()) if len(recientes) else 0

    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi1_valor"], f"{len(recientes)} caso" + ("s" if len(recientes) != 1 else ""))
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi1_label"], "Ingresados < 2 semanas")
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi2_valor"], f"{len(sin_gestion)} caso" + ("s" if len(sin_gestion) != 1 else ""))
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi2_label"], "Sin gestión inicial registrada")
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi3_valor"], f"~{dias_prom_recientes} días")
    _set_shape_text(s_gestiones, SLIDE_GESTIONES["kpi3_label"], "Antigüedad promedio del grupo")

    filas = []
    for i, (_, row) in enumerate(recientes.head(10).iterrows(), start=1):
        observacion = row["Observacion"] or row["Observacion_sugerida"] or "Sin gestión registrada"
        filas.append([
            i, row["Caso"], _fila_siniestro_nickname(row["Numero_Siniestro"], row["Nickname"]), row["Divisa"], row["Dias"],
            f"{int(row['Prob'])}%", _truncar_texto(observacion, 100),
        ])
    _fill_table_rows(_get_table(s_gestiones, SLIDE_GESTIONES["tabla"]), filas)

    # --- Slide 7: Próximos pasos y focos de gestión ---
    s7 = prs.slides[7]

    # --- Slides "Detalle Completo de Cartera" (una o más según la cantidad de
    # casos): se insertan justo antes de "Próximos pasos", clonando el
    # encabezado/KPIs de Tiempos de Residencia pero con una tabla propia más
    # densa (hasta 20 casos por página en vez de 10). ---
    s7_id = s7.slide_id
    detalle = tabla.sort_values("Dias", ascending=False).reset_index(drop=True)
    paginas_detalle = [
        detalle.iloc[i:i + _TABLA_DETALLE_FILAS_MAX] for i in range(0, len(detalle), _TABLA_DETALLE_FILAS_MAX)
    ]

    for num_pagina, pagina in enumerate(paginas_detalle, start=1):
        s_detalle = _duplicar_slide_antes_de(prs, 4, s7_id)
        _quitar_shape(s_detalle, SLIDE_GESTIONES["tabla"])
        _set_shape_text(s_detalle, SLIDE_GESTIONES["titulo"], "DETALLE COMPLETO DE CARTERA")
        _set_shape_text(
            s_detalle, SLIDE_GESTIONES["subtitulo"],
            f"Página {num_pagina} de {len(paginas_detalle)}  ·  {len(detalle)} casos",
        )
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi1_valor"], f"{len(detalle)} caso" + ("s" if len(detalle) != 1 else ""))
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi1_label"], "Total cartera filtrada")
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi2_valor"], f"{num_pagina} / {len(paginas_detalle)}")
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi2_label"], "Página")
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi3_valor"], f"{fmt_uf(kpis['uf_total'])}  +  {fmt_usd_m(kpis['usd_total'])}")
        _set_shape_text(s_detalle, SLIDE_GESTIONES["kpi3_label"], "Exposición total (UF + USD)")

        filas = []
        inicio = (num_pagina - 1) * _TABLA_DETALLE_FILAS_MAX + 1
        for i, (_, row) in enumerate(pagina.iterrows(), start=inicio):
            filas.append([
                i, row["Caso"], row["Numero_Siniestro"], row["Nickname"], row["Divisa"], row["Dias"], f"{int(row['Prob'])}%",
            ])
        _construir_tabla_detalle(s_detalle, filas)

    for i, ids in enumerate(SLIDE7["pasos"]):
        paso = pasos[i] if i < len(pasos) else {"titulo": "", "desc": ""}
        _set_shape_text(s7, ids["titulo"], paso.get("titulo", ""))
        _set_shape_text(s7, ids["desc"], paso.get("desc", ""))
    _set_shape_text(
        s7, SLIDE7["footer"],
        f"JPV Asociados · Ajustadores Especializados  ·  Análisis al {fmt_fecha_larga(fecha_corte)}  ·  Cartera {titulo_cartera}",
    )

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
